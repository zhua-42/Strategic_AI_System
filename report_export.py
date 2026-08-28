# -*- coding: utf-8 -*-
"""
报告文档导出引擎（Word / PPT / 图表渲染）· 2026 大改版
=====================================================
设计目标（参考 JPMC Forage《Task 2 JPMC.pptx》与券商研报版式）：

1. PPT 图文并茂：每张图表页 = 左侧图表 + 右侧简要点评文字（不是单页只有图）；
   底部统一「数据来源 Source」标注；所有表格/图片均限制在幻灯片边界内。
2. 配色：JPMC 风格——白底 + 深蓝主色 + 灰阶正文 + 强调色（沿用 Deloitte 绿/蓝系）。
3. Word：正文 markdown 表格自动转为真正的 Word 表格；清理 emoji/符号乱码；
   每张图表下方标注「图注 + 数据来源」。
4. 图表渲染：matplotlib（稳定、无需浏览器），新增龙头对比/杜邦瀑布/新闻来源等图型。

对外接口：
    render_chart_png(chart_type, data, ...) -> PNG bytes
    export_docx(...) -> Word bytes
    export_pptx(...) -> PPT bytes
"""
import io
import os
import re
import datetime

import pandas as pd

# ---------- 主题色（JPMC / Deloitte 融合） ----------
C_NAVY = "#0F2A5C"      # 深蓝（JPMC 主色，替代原来的 #1e3a8a）
C_BLUE = "#1F5FA8"      # 蓝
C_GREEN = "#86bc25"     # Deloitte 绿（强调）
C_TEAL = "#0d9488"      # 青
C_RED = "#C0392B"       # 红（警示/负值）
C_AMBER = "#E67E22"     # 琥珀
C_PURPLE = "#7D5BA6"    # 紫
C_SLATE = "#404040"     # 正文深灰（JPMC 用 #6D6E6A，加深以提升可读性）
C_MID = "#8A8F98"       # 次级灰
C_LIGHT = "#F4F6F9"     # 浅底
C_BORDER = "#C9CDD4"    # 表格边框
C_HEADER_BG = "#0F2A5C" # 表头底色

FONT_BUNDLED = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "assets", "fonts", "NotoSansSC-VF.ttf")
FONT_CANDIDATES = [
    FONT_BUNDLED,
    r"C:\Windows\Fonts\NotoSansSC-VF.ttf",
    r"C:\Windows\Fonts\msyh.ttc",
    r"C:\Windows\Fonts\simhei.ttf",
    "Noto Sans CJK SC", "Microsoft YaHei", "SimHei", "WenQuanYi Zen Hei",
]

_FONT_CACHE = {"name": None, "props": None}


def setup_font():
    """注册中文字体，返回 matplotlib FontProperties（缓存）。"""
    if _FONT_CACHE["props"] is not None:
        return _FONT_CACHE["props"]
    import matplotlib
    from matplotlib import font_manager as fm

    chosen = None
    for cand in FONT_CANDIDATES:
        try:
            if isinstance(cand, str) and os.path.isfile(cand):
                fm.fontManager.addfont(cand)
                name = fm.FontProperties(fname=cand).get_name()
                chosen = (cand, name)
                break
        except Exception:
            continue
    if chosen is None:
        for name in FONT_CANDIDATES:
            if not isinstance(name, str) or not os.path.exists(name):
                try:
                    fm.fontManager.findfont(name, fallback_to_default=False)
                    chosen = (None, name)
                    break
                except Exception:
                    continue
    if chosen is None:
        chosen = (None, "DejaVu Sans")
    matplotlib.rcParams["font.sans-serif"] = [chosen[1], "DejaVu Sans"]
    matplotlib.rcParams["axes.unicode_minus"] = False
    props = fm.FontProperties(fname=chosen[0]) if chosen[0] else fm.FontProperties(family=chosen[1])
    _FONT_CACHE["props"] = props
    return props


def _fig_png(fig, dpi=150):
    """matplotlib 图 → PNG bytes。"""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    return buf.getvalue()


def _style_ax(ax, title=None):
    ax.set_facecolor("white")
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(colors=C_SLATE, labelsize=9)
    ax.grid(axis="y", linestyle="--", alpha=0.35, color="#cbd5e1")
    if title:
        ax.set_title(title, fontsize=12, color=C_NAVY, fontweight="bold", pad=10)


def _fmt(v, nd=1):
    try:
        return f"{float(v):.{nd}f}"
    except Exception:
        return "—"


# ============================================================
# 图表渲染（matplotlib）
# 支持：dupont_compare / market_share / market_growth / financial_trend /
#       capability_compare / company_radar / risk_radar / industry_chain /
#       leader_compare / dupont_waterfall / news_source
# ============================================================
def render_chart_png(chart_type, data=None, is_company_mode=False, title="", subtitle=""):
    import matplotlib.pyplot as plt
    setup_font()
    data = data or {}
    fig = None

    def _make(figsize=(7.4, 4.2)):
        return plt.figure(figsize=figsize, dpi=150, facecolor="white")

    if chart_type == "dupont_compare":
        fig = _make()
        ax = fig.add_subplot(111)
        labels = ["ROE (%)", "净利润率 (%)", "资产周转率×100", "权益乘数×10"]
        comp = [data.get("company_roe", 0), data.get("company_margin", 0),
                data.get("company_turnover", 0) * 100, data.get("company_multiplier", 0) * 10]
        ind = [data.get("industry_roe", 0), data.get("industry_margin", 0),
               data.get("industry_turnover", 0) * 100, data.get("industry_multiplier", 0) * 10]
        x = range(len(labels))
        ax.bar([i - 0.19 for i in x], comp, width=0.38, color=C_NAVY, label=data.get("company_name", "标的公司"))
        ax.bar([i + 0.19 for i in x], ind, width=0.38, color=C_GREEN, label="行业均值基准")
        ax.set_xticks(list(x)); ax.set_xticklabels(labels, fontsize=9)
        _style_ax(ax, title or "标的公司与行业杜邦因子对比（标准化）")
        ax.legend(fontsize=9, loc="upper right", frameon=False)
        for i, v in enumerate(comp):
            ax.text(i - 0.19, v + 0.6, f"{v:.1f}", ha="center", fontsize=7.5, color=C_NAVY)
        for i, v in enumerate(ind):
            ax.text(i + 0.19, v + 0.6, f"{v:.1f}", ha="center", fontsize=7.5, color=C_GREEN)

    elif chart_type == "market_share":
        fig = _make()
        ax = fig.add_subplot(111)
        labels = data.get("labels", ["头部企业 (CR4)", "中坚力量", "尾部企业"])
        values = data.get("values", [50, 35, 15])
        colors = [C_GREEN, C_NAVY, "#cbd5e1"]
        wedges, _texts, autotexts = ax.pie(values, labels=labels, autopct="%1.1f%%",
                                           colors=colors, startangle=90, textprops={"fontsize": 9},
                                           wedgeprops=dict(width=0.42, edgecolor="white"))
        for at in autotexts:
            at.set_color("white"); at.set_fontweight("bold"); at.set_fontsize(9)
        ax.set_title(title or "行业市场集中度（CR4）动态格局", fontsize=12, color=C_NAVY, fontweight="bold", pad=10)

    elif chart_type == "market_growth":
        fig = _make()
        ax = fig.add_subplot(111)
        years = data.get("years", ["2022", "2023", "2024", "2025", "2026(E)"])
        size = data.get("market_size", [100, 110, 120, 130, 140])
        rate = data.get("growth_rate", [10, 10, 9, 8, 7])
        ax.bar(years, size, color=C_NAVY, alpha=0.85, label="市场规模 (亿元)")
        ax.set_ylabel("市场规模 (亿元)", color=C_NAVY, fontsize=9)
        ax2 = ax.twinx()
        ax2.plot(years, rate, color=C_RED, marker="o", linewidth=2.5, label="增速 (%)")
        ax2.set_ylabel("增速 (%)", color=C_RED, fontsize=9)
        ax2.spines["right"].set_visible(False); ax2.spines["top"].set_visible(False)
        ax2.tick_params(colors=C_RED, labelsize=9)
        _style_ax(ax, title or "行业市场规模与复合增速")
        ax.legend(fontsize=8, loc="upper left", frameon=False)
        ax2.legend(fontsize=8, loc="upper right", frameon=False)

    elif chart_type == "financial_trend":
        fig = _make()
        ax = fig.add_subplot(111)
        years = data.get("years", ["2022", "2023", "2024", "2025", "2026Q2"])
        roe = data.get("roe_trend", [12, 11, 10, 9.5, 9.1])
        margin = data.get("margin_trend", [10, 9.5, 9, 8.8, 8.5])
        ax.plot(years, roe, marker="o", color=C_BLUE, linewidth=2.5, label="平均 ROE (%)")
        ax.plot(years, margin, marker="s", color=C_TEAL, linewidth=2.5, label="净利润率 (%)")
        _style_ax(ax, title or "主要盈利指标变化趋势")
        ax.legend(fontsize=9, frameon=False)
        ax.set_ylim(bottom=0)

    elif chart_type == "capability_compare":
        fig = _make()
        ax = fig.add_subplot(111)
        metrics = data.get("metrics", ["盈利能力(ROE%)", "短期流动性(流动比率×10)", "资产效率(周转率×100)", "安全边际(现金流%)"])
        values = data.get("values", [12, 15, 60, 20])
        y = range(len(metrics))
        ax.barh(list(y), values, color=[C_GREEN, C_TEAL, C_BLUE, C_NAVY], height=0.55)
        ax.set_yticks(list(y)); ax.set_yticklabels(metrics, fontsize=9)
        for i, v in enumerate(values):
            ax.text(v + max(values) * 0.02, i, f"{v:.1f}", va="center", fontsize=8.5, color=C_SLATE)
        _style_ax(ax, title or "企业多维核心财务能力对比")
        ax.grid(axis="x", linestyle="--", alpha=0.3, color="#cbd5e1")
        ax.grid(axis="y", visible=False)

    elif chart_type == "company_radar":
        fig = _make((6.8, 4.6))
        ax = fig.add_subplot(111, polar=True)
        theta_vars = ["ROE", "净利润率", "资产周转率", "财务杠杆", "经营现金流"]
        comp_r = [data.get("company_roe", 0), data.get("company_margin", 0),
                  data.get("company_turnover", 0) * 10, data.get("company_multiplier", 0),
                  data.get("company_cash", 0) / 1000]
        ind_r = [data.get("industry_roe", 0), data.get("industry_margin", 0),
                 data.get("industry_turnover", 0) * 10, data.get("industry_multiplier", 0),
                 data.get("industry_cash", 0) / 1000]
        angles = [i / len(theta_vars) * 2 * 3.1415926 for i in range(len(theta_vars))]
        angles_closed = angles + angles[:1]
        comp_c = comp_r + comp_r[:1]
        ind_c = ind_r + ind_r[:1]
        ax.plot(angles_closed, comp_c, color=C_NAVY, linewidth=2, label=data.get("company_name", "标的公司"))
        ax.fill(angles_closed, comp_c, color=C_NAVY, alpha=0.15)
        ax.plot(angles_closed, ind_c, color=C_GREEN, linewidth=2, label="行业平均")
        ax.fill(angles_closed, ind_c, color=C_GREEN, alpha=0.12)
        ax.set_xticks(angles); ax.set_xticklabels(theta_vars, fontsize=9, color=C_SLATE)
        ax.set_ylim(0, max(50.0, max(comp_r + ind_r + [1]) * 1.4))
        ax.set_title(title or "标的公司与行业能力多维透视", fontsize=12, color=C_NAVY, fontweight="bold", pad=18)
        ax.legend(loc="upper right", bbox_to_anchor=(1.32, 1.1), fontsize=8.5, frameon=False)

    elif chart_type == "risk_radar":
        fig = _make((7.0, 4.6))
        ax = fig.add_subplot(111, polar=True)
        dims = data.get("dimensions", ["偿债与财务杠杆风险", "短期流动性紧缺风险", "存货/资产减值风险",
                                       "盈利质量恶化风险", "政策合规与壁垒风险"])
        vals = data.get("values", [3, 3, 3, 3, 3])
        angles = [i / len(dims) * 2 * 3.1415926 for i in range(len(dims))]
        angles_closed = angles + angles[:1]
        vals_closed = vals + vals[:1]
        ax.plot(angles_closed, vals_closed, color=C_RED, linewidth=2.2)
        ax.fill(angles_closed, vals_closed, color=C_RED, alpha=0.25)
        ax.set_xticks(angles); ax.set_xticklabels(dims, fontsize=8.6, color=C_SLATE)
        ax.set_ylim(0, 5.0)
        ax.set_yticks([1, 2, 3, 4, 5])
        ax.set_yticklabels(["1 安全", "2", "3 关注", "4", "5 高危"], fontsize=7.5, color="#94a3b8")
        ax.set_title(title or "企业经营及财务多维风险指数（1=安全，5=高危）",
                     fontsize=12, color=C_RED, fontweight="bold", pad=18)
        ax.grid(color="#e2e8f0", alpha=0.8)

    elif chart_type == "industry_chain":
        fig = _make((8.6, 3.8))
        ax = fig.add_subplot(111)
        nodes = data.get("nodes", [])
        if not nodes:
            ax.text(0.5, 0.5, "该行业暂无产业链明细数据", ha="center", va="center", fontsize=11, color=C_SLATE)
            ax.axis("off")
        else:
            ax.axis("off")
            stage_color = {
                "upstream": C_BLUE, "midstream": C_TEAL,
                "integration": C_AMBER, "downstream": C_NAVY, "service": C_PURPLE,
            }
            for i, nd in enumerate(nodes):
                x = i / max(1, len(nodes) - 1)
                color = stage_color.get(nd.get("stage", ""), C_BLUE)
                ax.scatter([x], [0.5], s=1250, color=color, alpha=0.92, edgecolors="white", zorder=3)
                ax.text(x, 0.5, f"{nd.get('name', '')[:8]}\n{nd.get('name', '')[8:16]}",
                        ha="center", va="center", fontsize=8.2, color="white", fontweight="bold", zorder=4)
                ax.text(x, 0.72, f"龙头: {nd.get('leaders', '')[:26]}",
                        ha="center", va="center", fontsize=7.0, color=C_SLATE, zorder=4)
                ax.text(x, 0.24, f"利润率: {nd.get('margin', '')[:22]}",
                        ha="center", va="center", fontsize=7.0, color=C_SLATE, zorder=4)
                if i < len(nodes) - 1:
                    nx = (i + 1) / max(1, len(nodes) - 1)
                    ax.annotate("", xy=(nx, 0.5), xytext=(x + 0.04, 0.5),
                                arrowprops=dict(arrowstyle="-|>", color="#94a3b8", lw=2))
            ax.set_xlim(-0.03, 1.03); ax.set_ylim(0.1, 0.9)
            ax.set_title(title or "产业链全景逻辑流（环节→龙头→利润率）",
                         fontsize=12, color=C_NAVY, fontweight="bold", pad=12)

    elif chart_type == "leader_compare":
        """龙头公司横向对比：多指标分组柱状图。"""
        fig = _make((8.0, 4.4))
        ax = fig.add_subplot(111)
        companies = data.get("companies", [])
        metrics = data.get("metrics", [])
        values = data.get("values", {})
        if companies and metrics:
            x = range(len(metrics))
            width = 0.8 / max(1, len(companies))
            colors = [C_NAVY, C_BLUE, C_TEAL, C_GREEN, C_AMBER, C_PURPLE]
            for ci, comp in enumerate(companies):
                ys = [values.get(comp, {}).get(m, 0) for m in metrics]
                ax.bar([i + (ci - (len(companies) - 1) / 2) * width for i in x],
                       ys, width=width * 0.9, label=comp, color=colors[ci % len(colors)])
            ax.set_xticks(list(x)); ax.set_xticklabels(metrics, fontsize=8.5)
            ax.legend(fontsize=8, loc="upper right", frameon=False, ncol=2)
            ax.set_ylabel("指标值", fontsize=9)
        _style_ax(ax, title or "龙头公司横向对比（ROE / 毛利率 / 净利率 / 同比 / EPS）")
        ax.grid(axis="y", linestyle="--", alpha=0.3, color="#cbd5e1")

    elif chart_type == "dupont_waterfall":
        """杜邦拆解瀑布图（公司 vs 行业 ROE 差距归因）。"""
        fig = _make((8.0, 4.2))
        ax = fig.add_subplot(111)
        comp_roe = data.get("company_roe", 0)
        ind_roe = data.get("industry_roe", 0)
        comp_m = data.get("company_margin", 0)
        ind_m = data.get("industry_margin", 0)
        comp_t = data.get("company_turnover", 0)
        ind_t = data.get("industry_turnover", 0)
        comp_l = data.get("company_multiplier", 0)
        ind_l = data.get("industry_multiplier", 0)
        # 相对行业基准的贡献分解（对数近似：ln 差 = 各因子 ln 差之和）
        import math
        def _safe_ln(v):
            return math.log(max(v, 0.01))
        try:
            total = _safe_ln(comp_roe) - _safe_ln(ind_roe) if comp_roe > 0 and ind_roe > 0 else 0
            d_m = _safe_ln(comp_m) - _safe_ln(ind_m) if comp_m > 0 and ind_m > 0 else 0
            d_t = _safe_ln(comp_t) - _safe_ln(ind_t) if comp_t > 0 and ind_t > 0 else 0
            d_l = _safe_ln(comp_l) - _safe_ln(ind_l) if comp_l > 0 and ind_l > 0 else 0
        except Exception:
            d_m = d_t = d_l = 0
            total = 0
        labels = ["行业ROE", "利润率贡献", "周转率贡献", "杠杆贡献", "公司ROE"]
        vals = [ind_roe, d_m, d_t, d_l, comp_roe]
        # 瀑布：累计
        cum = ind_roe
        bottoms = [0, ind_roe, ind_roe + d_m, ind_roe + d_m + d_t, 0]
        heights = [ind_roe, d_m, d_t, d_l, comp_roe]
        for i, (lb, v) in enumerate(zip(labels, vals)):
            if i in (1, 2, 3):
                color = C_GREEN if v >= 0 else C_RED
                ax.bar(i, abs(v), bottom=bottoms[i], color=color, alpha=0.85, width=0.55)
                ax.text(i, bottoms[i] + abs(v) / 2, f"{v:+.2f}", ha="center", va="center",
                        fontsize=9, color="white", fontweight="bold")
            else:
                ax.bar(i, abs(v), bottom=0, color=C_NAVY if i == 0 else C_BLUE, width=0.55)
                ax.text(i, abs(v) + 0.4, f"{v:.2f}", ha="center", fontsize=9, color=C_NAVY, fontweight="bold")
        ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels, fontsize=9)
        _style_ax(ax, title or "杜邦 ROE 差距归因（公司 vs 行业，对数分解示意）")
        ax.grid(axis="y", linestyle="--", alpha=0.3, color="#cbd5e1")

    elif chart_type == "news_source":
        """新闻/公告来源分布（环形图）。"""
        fig = _make((6.4, 3.6))
        ax = fig.add_subplot(111)
        sources = data.get("sources", {}) or {}
        if sources:
            labels = list(sources.keys())
            values = list(sources.values())
            colors = [C_NAVY, C_BLUE, C_TEAL, C_GREEN, C_AMBER, C_PURPLE, "#94a3b8"][:len(labels)]
            ax.pie(values, labels=labels, autopct="%1.0f%%", colors=colors, startangle=90,
                   textprops={"fontsize": 8.5}, wedgeprops=dict(width=0.45, edgecolor="white"))
            ax.set_title(title or "新闻/公告信息来源分布", fontsize=12, color=C_NAVY, fontweight="bold", pad=10)
        else:
            ax.text(0.5, 0.5, "暂无新闻数据", ha="center", va="center", fontsize=11, color=C_SLATE)
            ax.axis("off")

    if fig is None:
        fig = _make()
        ax = fig.add_subplot(111)
        ax.text(0.5, 0.5, "暂无图表数据", ha="center", va="center", fontsize=12, color=C_SLATE)
        ax.axis("off")
    png = _fig_png(fig)
    import matplotlib.pyplot as plt
    plt.close(fig)
    return png


# ---------- 高级导出：把 Plotly 图转 PNG（优先 kaleido，失败回退） ----------
def fig_plotly_to_png(fig, width=900, height=560):
    """优先 kaleido 渲染（云端 0.2.1 可用），失败返回 None。"""
    try:
        return fig.to_image(format="png", width=width, height=height)
    except Exception:
        return None


def chart_png_or_none(fig, chart_type, data=None, fallback_title=""):
    """统一入口：kaleido → matplotlib 兜底，保证永远有图。"""
    png = fig_plotly_to_png(fig)
    if png:
        return png
    try:
        return render_chart_png(chart_type, data, title=fallback_title)
    except Exception:
        return None


# ---------- 文本清理（解决导出乱码/符号问题） ----------
_EMOJI_RE = re.compile(
    "[\U0001F300-\U0001FAFF\u2600-\u27BF\uFE0F\u2B00-\u2BFF\u2190-\u21FF\u2700-\u27BF]"
)


def clean_export_text(text, keep_markdown_tables=False):
    """
    清理导出文本中的乱码符号：
    - 去掉 emoji 与装饰符号
    - 去掉 markdown 标记（# * ` > --- 等）
    - 折叠多余空行/空白
    keep_markdown_tables=True 时保留 markdown 表格行（由调用方转 Word 表格）。
    """
    if not text:
        return ""
    t = str(text)
    t = _EMOJI_RE.sub("", t)
    t = re.sub(r"https?://\S+", "", t)  # 正文里去掉裸链接（表格里有来源列）
    if not keep_markdown_tables:
        # 把 markdown 表格行合并为一行文本
        lines = []
        for ln in t.split("\n"):
            s = ln.strip()
            if s.startswith("|") and s.endswith("|"):
                cells = [c.strip() for c in s.strip("|").split("|")]
                if all(re.fullmatch(r":?-{2,}:?", c or "-") for c in cells):
                    continue  # 分隔行
                lines.append(" | ".join(cells))
            else:
                lines.append(s)
        t = "\n".join(lines)
    t = re.sub(r"^#{1,6}\s*", "", t, flags=re.M)          # 标题井号
    t = re.sub(r"\*\*([^*]+)\*\*", r"\1", t)              # 加粗
    t = re.sub(r"\*([^*]+)\*", r"\1", t)                  # 斜体
    t = re.sub(r"`([^`]+)`", r"\1", t)                    # 行内代码
    t = re.sub(r"^>\s?", "", t, flags=re.M)               # 引用
    t = re.sub(r"^\s*[-*+]\s+", "• ", t, flags=re.M)      # 无序列表
    t = re.sub(r"^\s*\d+\.\s+", "", t, flags=re.M)        # 有序列表序号
    t = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", t)       # 图片
    t = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", t)        # 链接
    t = re.sub(r"^\s*[-=]{3,}\s*$", "", t, flags=re.M)    # 分隔线
    t = re.sub(r"[ \t]{2,}", " ", t)                      # 多余空格
    t = re.sub(r"\n{3,}", "\n\n", t)                      # 多余空行
    return t.strip()


def split_markdown_tables(report_text):
    """
    从报告中抽出 markdown 表格，返回 (clean_text, tables)。
    tables: [{"title", "headers", "rows"}]
    """
    tables = []
    lines = (report_text or "").split("\n")
    out_lines = []
    i = 0
    cur_title = ""
    while i < len(lines):
        ln = lines[i].strip()
        if ln.startswith("|") and ln.endswith("|"):
            # 收集连续表格行
            block = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                block.append(lines[i].strip())
                i += 1
            parsed = _parse_md_table(block)
            if parsed:
                parsed["title"] = cur_title
                tables.append(parsed)
            continue
        if re.match(r"^#{1,6}\s+", ln):
            cur_title = re.sub(r"^#{1,6}\s*", "", ln).strip()
        out_lines.append(lines[i])
        i += 1
    return "\n".join(out_lines), tables


def _parse_md_table(block):
    """把 markdown 表格块解析为 {headers, rows}。"""
    rows = []
    for ln in block:
        cells = [c.strip() for c in ln.strip().strip("|").split("|")]
        rows.append(cells)
    if not rows:
        return None
    # 去掉分隔行 |---|
    rows = [r for r in rows if not all(re.fullmatch(r":?-{2,}:?", c or "-") for c in r)]
    if not rows:
        return None
    headers = rows[0]
    body = rows[1:]
    return {"headers": headers, "rows": body}


def _excerpt(text, limit=180):
    t = re.sub(r"\s+", " ", text).strip()
    return t[:limit] + ("…" if len(t) > limit else "")


def _today():
    return datetime.date.today().strftime("%Y-%m-%d")


# ---------------- Word ----------------
def export_docx(query, report_text, chart_images, evidence_data=None, gap_data=None,
                meta=None, is_company_mode=True, source_text=""):
    """
    chart_images: dict -> {"chart_id": {"title","caption","png","source"}}
    升级：
    - markdown 表格自动转 Word 表格
    - 清理 emoji/符号
    - 每图附「图注 + 数据来源」
    - 证据链/缺口用真实表格呈现
    """
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn

    doc = Document()
    for section in doc.sections:
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)
    style = doc.styles["Normal"]
    style.font.name = "Microsoft YaHei"
    style.element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
    style.font.size = Pt(10.5)

    # ---- 封面 ----
    for _ in range(3):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"{query}\n深度战略研报")
    run.font.size = Pt(26); run.font.color.rgb = RGBColor(0x0F, 0x2A, 0x5C); run.bold = True
    doc.add_paragraph()
    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run(f"多智能体智能投研系统 · 数据可溯源版\n生成日期：{_today()}")
    r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0x86, 0xBC, 0x25)
    doc.add_paragraph()
    doc.add_paragraph("—" * 40)

    # ---- 数据口径与来源 ----
    doc.add_heading("数据口径与来源说明", level=1)
    if isinstance(meta, dict):
        for k, v in meta.items():
            doc.add_paragraph(f"• {k}：{v}", style=None)
    else:
        doc.add_paragraph(str(meta or ""))

    # ---- 数据看板 ----
    doc.add_heading("第一部分：数据看板可视化", level=1)
    if chart_images:
        for cid, ch in chart_images.items():
            doc.add_heading(clean_export_text(ch.get("title", cid)), level=2)
            if ch.get("png"):
                _add_chart_block(doc, ch["png"], ch.get("caption", ""),
                                 ch.get("source", ""))
            else:
                doc.add_paragraph("[图表数据缺失]")
    else:
        doc.add_paragraph("本次研究未生成图表。")

    # ---- 证据链（真实表格）----
    doc.add_heading("第二部分：数据可信度证据链（Evidence Ledger）", level=1)
    if evidence_data:
        table = doc.add_table(rows=1, cols=4)
        table.style = "Light Grid Accent 3"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        hdr = table.rows[0].cells
        for i, t in enumerate(["审计论点", "数据来源", "位置/页码", "可信度评级"]):
            hdr[i].text = t
            for para in hdr[i].paragraphs:
                for rr in para.runs:
                    rr.font.bold = True
        for item in evidence_data:
            cells = table.add_row().cells
            cells[0].text = clean_export_text(str(item.get("point", "")))
            cells[1].text = clean_export_text(str(item.get("source", "")))
            cells[2].text = clean_export_text(str(item.get("page", "")))
            cells[3].text = str(item.get("confidence", ""))
    else:
        doc.add_paragraph("暂无证据链条目。")

    # ---- 正文（markdown 表格转 Word 表格）----
    doc.add_heading("第三部分：深度研究报告正文", level=1)
    clean_body, md_tables = split_markdown_tables(report_text)
    # 段落级渲染
    sections = _split_sections(clean_body)
    for lvl, title, body in sections:
        doc.add_heading(f"{clean_export_text(title)}", level=min(lvl + 1, 4))
        for line in body:
            if line.strip():
                doc.add_paragraph(clean_export_text(line))
    # 报告内嵌表格
    for tbl in md_tables:
        if tbl.get("headers") and tbl.get("rows"):
            t = doc.add_table(rows=1, cols=len(tbl["headers"]))
            t.style = "Light Grid Accent 3"
            hdr = t.rows[0].cells
            for j, h in enumerate(tbl["headers"]):
                hdr[j].text = clean_export_text(h)
                for para in hdr[j].paragraphs:
                    for rr in para.runs:
                        rr.font.bold = True
            for row in tbl["rows"]:
                cells = t.add_row().cells
                for j in range(len(tbl["headers"])):
                    v = row[j] if j < len(row) else ""
                    cells[j].text = clean_export_text(v)
            doc.add_paragraph()

    # ---- 资料缺口 ----
    doc.add_heading("第四部分：资料缺口说明", level=1)
    if gap_data:
        for g in gap_data:
            doc.add_paragraph(clean_export_text(f"• {g}"))
    else:
        doc.add_paragraph("本次研究资料完整，无关键缺口。")

    # ---- 免责声明 ----
    doc.add_heading("免责声明", level=2)
    doc.add_paragraph("本报告由 AI 多智能体系统自动生成，所有数据来源已在正文与证据链中标注；"
                      "内容仅供学习与研究参考，不构成任何投资建议。投资决策与风险由使用者自行承担。")

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


# ---------------- PPT（JPMC 风格） ----------------
def export_pptx(query, report_text, chart_images, evidence_data=None, gap_data=None,
                meta=None, is_company_mode=True, source_text="", leader_data=None,
                news_items=None):
    """
    JPMC 风格 16:9 PPT：
    - 封面/目录/核心指标/逐图页（左图右文）/证据链/缺口/新闻/免责
    - 每页底部统一「数据来源 Source」行
    - 图表页：左侧图片 + 右侧要点文字（图文并茂）
    - 表格使用真实 PPT 表格，自动收缩字号避免超出边界
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    SW, SH = Inches(13.333), Inches(7.5)
    MARGIN = Inches(0.55)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    def _rgb(hexstr):
        h = hexstr.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _set_font(run, size=14, bold=False, color=C_SLATE, name="Microsoft YaHei"):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        run.font.name = name
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", name)

    def _bg(slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb("#FFFFFF")
        shape.line.fill.background()
        shape.shadow.inherit = False
        # 顶部 JPMC 风格深蓝细条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.09))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(C_NAVY)
        bar.line.fill.background()
        bar.shadow.inherit = False
        # 左下角品牌
        brand = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(4.0), Inches(0.32))
        bf = brand.text_frame
        bf.word_wrap = True
        bp = bf.paragraphs[0]
        br = bp.add_run()
        br.text = "数智投研多智能体系统 · AI Research"
        _set_font(br, size=8.5, color=C_MID)

    def _title(slide, text, sub=None, page=None):
        tb = slide.shapes.add_textbox(MARGIN, Inches(0.32), Inches(11.5), Inches(0.62))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = clean_export_text(text)
        _set_font(r, size=24, bold=True, color=C_NAVY)
        if sub:
            tb2 = slide.shapes.add_textbox(MARGIN + Inches(0.06), Inches(0.98), Inches(11.0), Inches(0.42))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = clean_export_text(sub)
            _set_font(r2, size=12.5, color=C_MID)
        if page:
            pgn = slide.shapes.add_textbox(Inches(12.35), Inches(0.36), Inches(0.7), Inches(0.4))
            pgf = pgn.text_frame
            pgr = pgf.paragraphs[0].add_run()
            pgr.text = str(page)
            _set_font(pgr, size=11, bold=True, color=C_MID)

    def _source_footer(slide, source_text):
        tb = slide.shapes.add_textbox(MARGIN, Inches(6.75), Inches(12.2), Inches(0.34))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "数据来源 Source: " + clean_export_text(source_text or "东方财富业绩报表 / 公司财报 / 公开资料")
        _set_font(r, size=9, color=C_MID)

    def _add_textbox(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
        """lines: [{"text","size","bold","color","bullet"}]"""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            r = p.add_run()
            text = ln.get("text", "")
            if ln.get("bullet") and text:
                text = "▪ " + text
            r.text = clean_export_text(text)
            _set_font(r, size=ln.get("size", 12), bold=ln.get("bold", False),
                      color=ln.get("color", C_SLATE))
        return tb

    def _fit_picture(slide, png, left, top, max_w_in, max_h_in):
        """按比例缩放图片，保证不超出 max 边界。"""
        from PIL import Image
        try:
            img = Image.open(io.BytesIO(png))
            w, h = img.size
            ratio = h / w if w else 1
            max_w, max_h = Inches(max_w_in), Inches(max_h_in)
            if ratio > max_h / max_w:
                pic_h = int(max_h)
                pic_w = int(max_h / ratio)
            else:
                pic_w = int(max_w)
                pic_h = int(max_w * ratio)
            slide.shapes.add_picture(io.BytesIO(png), left, top, pic_w, pic_h)
        except Exception:
            try:
                slide.shapes.add_picture(io.BytesIO(png), left, top, Inches(max_w_in), Inches(max_h_in))
            except Exception:
                pass

    def _add_table(slide, headers, rows, left, top, width_in, font=9.5, header_font=10,
                   row_h_in=0.3, max_rows=12):
        """添加真实 PPT 表格，自动收缩字号防溢出；返回 (table, shown_rows)。"""
        from pptx.util import Inches as _In
        rows = rows[:max_rows]
        n_rows = len(rows) + 1
        n_cols = max(len(headers), max((len(r) for r in rows), default=0))
        if n_cols == 0:
            return None, 0
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, _In(width_in), _In(row_h_in * n_rows))
        table = table_shape.table
        # 列宽均分
        col_w = _In(width_in) / n_cols
        for c in range(n_cols):
            table.columns[c].width = int(col_w)
        # 表头
        for c in range(n_cols):
            cell = table.cell(0, c)
            cell.text = clean_export_text(headers[c] if c < len(headers) else "")
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(C_HEADER_BG)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    _set_font(run, size=header_font, bold=True, color="#FFFFFF")
        # 数据行
        for ri, row in enumerate(rows):
            fill = _rgb(C_LIGHT) if ri % 2 == 1 else _rgb("#FFFFFF")
            for c in range(n_cols):
                cell = table.cell(ri + 1, c)
                cell.text = clean_export_text(row[c] if c < len(row) else "")
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        _set_font(run, size=font, color=C_SLATE)
        return table, len(rows)

    # ================= 1) 封面 =================
    s = prs.slides.add_slide(blank)
    _bg(s)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SW, Inches(2.4))
    band.fill.solid(); band.fill.fore_color.rgb = _rgb(C_NAVY)
    band.line.fill.background(); band.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.75), Inches(11.3), Inches(1.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = clean_export_text(f"{query}")
    _set_font(r, size=40, bold=True, color="#FFFFFF")
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = "深度战略研报 · Deep Research Report"
    _set_font(r2, size=18, color="#C9D4E8")
    tb2 = s.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.9))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    r3 = p3.add_run()
    r3.text = f"多智能体智能投研系统 · 数据可溯源版    生成日期：{_today()}"
    _set_font(r3, size=14, color=C_MID)
    _source_footer(s, "AI 生成内容，仅供学习与研究参考，不构成投资建议")

    # ================= 2) 目录 =================
    s = prs.slides.add_slide(blank)
    _bg(s); _title(s, "目录 · Contents", page=2)
    items = ["01 数据口径与来源", "02 核心指标总览", "03 数据看板可视化（图表 + 要点）",
             "04 龙头公司横向对比", "05 最新新闻与公告", "06 深度研究报告正文",
             "07 数据可信度证据链", "08 资料缺口说明", "09 免责声明"]
    lines = [{"text": it, "size": 17, "color": C_NAVY if i % 2 == 0 else C_SLATE, "bullet": True}
             for i, it in enumerate(items)]
    _add_textbox(s, Inches(1.2), Inches(1.6), Inches(11.0), Inches(4.8), lines)
    _source_footer(s, "数智投研多智能体系统 · 生成日期 " + _today())

    # ================= 3) 数据口径 =================
    s = prs.slides.add_slide(blank)
    _bg(s); _title(s, "01 数据口径与来源", page=3)
    meta_lines = []
    if isinstance(meta, dict):
        for k, v in meta.items():
            meta_lines.append({"text": f"{k}：{v}", "size": 13, "color": C_SLATE, "bullet": True})
    meta_lines.append({"text": "风险维度采用真实财务指标映射（权益乘数/现金流质量/毛利率/ROE/政策基准）",
                       "size": 12, "color": C_MID, "bullet": True})
    _add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.6), meta_lines)
    _source_footer(s, "东方财富业绩报表聚合 / 本地 SQLite / RAG 知识库")

    # ================= 4) 图表页（左图右文，图文并茂） =================
    page_no = 4
    if chart_images:
        for cid, ch in chart_images.items():
            s = prs.slides.add_slide(blank)
            _bg(s)
            _title(s, ch.get("title", cid), ch.get("caption", ""), page=page_no)
            if ch.get("png"):
                # 左图：最大 8.0in 宽 x 4.6in 高（留出右边文字与底部来源）
                _fit_picture(s, ch["png"], Inches(0.55), Inches(1.75), 7.9, 4.55)
            # 右侧要点框
            notes = ch.get("notes") or _auto_notes(cid, ch)
            if notes:
                panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(8.7), Inches(1.75), Inches(4.05), Inches(4.55))
                panel.fill.solid(); panel.fill.fore_color.rgb = _rgb(C_LIGHT)
                panel.line.color.rgb = _rgb(C_BORDER); panel.line.width = Pt(1)
                panel.shadow.inherit = False
                tfb = panel.text_frame
                tfb.word_wrap = True
                tfb.margin_left = Inches(0.15); tfb.margin_right = Inches(0.15)
                tfb.margin_top = Inches(0.12); tfb.margin_bottom = Inches(0.12)
                for i, nt in enumerate(notes[:7]):
                    p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
                    p.space_after = Pt(5)
                    rr = p.add_run()
                    rr.text = "▪ " + clean_export_text(nt)
                    _set_font(rr, size=10.5, color=C_SLATE)
            _source_footer(s, ch.get("source") or "数据来源见各图口径说明")
            page_no += 1

    # ================= 5) 龙头对比 =================
    if leader_data and leader_data.get("companies"):
        s = prs.slides.add_slide(blank)
        _bg(s)
        _title(s, "04 龙头公司横向对比", f"行业：{leader_data.get('industry', '')} · 3~4 家龙头",
               page=page_no)
        comps = leader_data["companies"]
        metrics = leader_data.get("metrics", [])
        values = leader_data.get("values", {})
        if comps and metrics:
            rows = []
            for m in metrics:
                rows.append([m] + [_fmt(values.get(c, {}).get(m), 1) for c in comps])
            headers = ["指标"] + list(comps)
            _add_table(s, headers, rows, Inches(0.55), Inches(1.7), 12.2, font=9.5, row_h_in=0.32)
        notes = [f"{c}：{leader_data.get('notes', {}).get(c, '')}" for c in comps]
        if notes:
            lines = [{"text": n, "size": 10.5, "color": C_SLATE, "bullet": True} for n in notes[:6]]
            _add_textbox(s, Inches(0.55), Inches(5.3), Inches(12.2), Inches(1.3), lines)
        _source_footer(s, "；".join(leader_data.get("sources", [])) or "东方财富财务摘要接口 / 兜底基准")
        page_no += 1

    # ================= 6) 新闻与公告 =================
    if news_items:
        s = prs.slides.add_slide(blank)
        _bg(s)
        _title(s, "05 最新新闻与公告", "实时公开信息 · 东方财富 / 搜狗 / 财新", page=page_no)
        rows = []
        for it in news_items[:10]:
            rows.append([str(it.get("date", ""))[:10], str(it.get("title", ""))[:36],
                         str(it.get("source", ""))[:14]])
        if rows:
            _add_table(s, ["日期", "标题", "来源"], rows, Inches(0.55), Inches(1.7), 12.2,
                       font=9.5, row_h_in=0.34, max_rows=11)
        _source_footer(s, "东方财富公告大全 / 7×24 快讯 / 搜狗新闻 / 财新网")
        page_no += 1

    # ================= 7) 正文摘要 =================
    sections = _split_sections(report_text)
    s = prs.slides.add_slide(blank)
    _bg(s); _title(s, "06 深度研究报告正文（摘要）", page=page_no)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.9))
    tf = tb.text_frame; tf.word_wrap = True
    count = 0
    for lvl, title, body in sections[:9]:
        full = clean_export_text(" ".join(body))
        if not full.strip():
            continue
        p = tf.paragraphs[0] if count == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run(); r.text = f"▎{clean_export_text(title)}"
        _set_font(r, size=14, bold=True, color=C_NAVY)
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = _excerpt(full, 160)
        _set_font(r2, size=11, color=C_SLATE)
        count += 1
    _source_footer(s, f"{query} · {_today()}")
    page_no += 1

    # ================= 8) 证据链（表格） =================
    s = prs.slides.add_slide(blank)
    _bg(s); _title(s, "07 数据可信度证据链", page=page_no)
    if evidence_data:
        rows = []
        for item in evidence_data[:10]:
            rows.append([str(item.get("point", ""))[:30], str(item.get("source", ""))[:18],
                         str(item.get("page", ""))[:14], str(item.get("confidence", ""))])
        _add_table(s, ["审计论点", "数据来源", "位置/页码", "可信度"], rows,
                   Inches(0.55), Inches(1.7), 12.2, font=9.5, row_h_in=0.34, max_rows=11)
    else:
        _add_textbox(s, Inches(0.8), Inches(2.0), Inches(11.8), Inches(1.0),
                     [{"text": "暂无证据链条目", "size": 14, "color": C_SLATE}])
    _source_footer(s, "专家委员会终审 · A=数据库强锁定 B=RAG校验 C=模型推断")
    page_no += 1

    # ================= 9) 资料缺口 =================
    s = prs.slides.add_slide(blank)
    _bg(s); _title(s, "08 资料缺口说明", page=page_no)
    if gap_data:
        rows = []
        for g in gap_data:
            item = str(g)
            if "（" in item:
                name, rest = item.split("（", 1)
                rows.append([name[:28], rest.rstrip("）")[:40]])
            else:
                rows.append([item[:28], ""])
        _add_table(s, ["缺口项", "原因/建议"], rows, Inches(0.55), Inches(1.7), 12.2,
                   font=9.5, row_h_in=0.34, max_rows=11)
    else:
        _add_textbox(s, Inches(0.8), Inches(2.0), Inches(11.8), Inches(1.0),
                     [{"text": "本次研究资料完整，无关键缺口。", "size": 14, "color": C_SLATE}])
    _source_footer(s, "资料缺口审查 · 数据源限制如实标注")
    page_no += 1

    # ================= 10) 免责声明 =================
    s = prs.slides.add_slide(blank)
    _bg(s); _title(s, "09 免责声明", page=page_no)
    _add_textbox(s, Inches(0.8), Inches(2.0), Inches(11.8), Inches(2.6), [
        {"text": "本报告由 AI 多智能体系统自动生成，所有数据来源已在正文、图表与证据链中标注；",
         "size": 15, "color": C_SLATE},
        {"text": "内容仅供学习与研究参考，不构成任何投资建议。投资决策与风险由使用者自行承担。",
         "size": 15, "color": C_SLATE},
        {"text": "数据截至报告期：见各图表来源标注。", "size": 12, "color": C_MID},
    ])
    _source_footer(s, "AI 生成 · 仅供参考")

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _auto_notes(cid, ch):
    """图表要点自动生成（图文并茂的右侧文字）。"""
    caption = str(ch.get("caption", "") or "")
    source = str(ch.get("source", "") or "")
    notes = []
    if caption:
        notes.append(caption)
    notes.append("上图展示了该维度的核心数据与对比口径。")
    if source:
        notes.append("数据来源：" + source)
    return notes[:5]


def _add_chart_block(doc_or_slide, png, caption, source=""):
    from docx.shared import Inches
    from io import BytesIO
    try:
        doc_or_slide.add_picture(BytesIO(png), width=Inches(5.9))
        cap = f"图注：{caption}"
        if source:
            cap += f"\n数据来源：{source}"
        doc_or_slide.add_paragraph(cap)
    except Exception as e:
        doc_or_slide.add_paragraph(f"[图表插入失败: {e}]")


def _split_sections(report_text):
    """把 markdown 报告切成 (level, title, body_lines)。"""
    sections = []
    cur_title = "报告正文"
    cur_level = 2
    body = []
    for line in (report_text or "").split("\n"):
        s = line.strip()
        if re.match(r"^#+\s+", s):
            if body and any(b.strip() for b in body):
                sections.append((cur_level, cur_title, body))
            lvl = len(s) - len(s.lstrip("#"))
            cur_title = re.sub(r"^#+\s*", "", s).strip()
            cur_level = lvl
            body = []
        else:
            body.append(s)
    if body and any(b.strip() for b in body):
        sections.append((cur_level, cur_title, body))
    return sections or [(2, "报告正文", (report_text or "").split("\n"))]


if __name__ == "__main__":
    # 自检
    png = render_chart_png("risk_radar", {"dimensions": ["a", "b", "c", "d", "e"], "values": [3, 4, 2, 3.5, 4]})
    print("risk_radar png:", len(png) if png else None)
    png2 = render_chart_png("dupont_compare", {"company_roe": 18.5, "company_margin": 5.2, "company_turnover": 1.1,
                                               "company_multiplier": 2.5, "industry_roe": 12.5, "industry_margin": 8.2,
                                               "industry_turnover": 0.75, "industry_multiplier": 2.1, "company_name": "比亚迪"})
    print("dupont png:", len(png2) if png2 else None)
    # 龙头对比图
    try:
        import leader_compare
        lp = leader_compare.build_leader_payload("新能源汽车")
        png3 = render_chart_png("leader_compare", lp, title="龙头公司横向对比")
        print("leader png:", len(png3) if png3 else None)
    except Exception as e:
        print("leader png skipped:", e)
